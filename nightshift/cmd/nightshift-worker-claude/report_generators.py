"""Helpers for generating PDF / DOCX / XLSX report artifacts.

Each entry point returns `(binary_bytes, preview_html)` so the caller can
upload the binary as the authoritative artifact and the HTML as a
companion inline preview (stored next to the binary in S3).

The HTML preview is deliberately the same shape for PDF and DOCX — both
go through the same markdown -> HTML pipeline, and both use a shared
print-friendly wrapper so the inline preview looks like the document
the user will open.

XLSX uses its own preview format: a self-contained tabbed HTML document
with one <table> per sheet.
"""

from __future__ import annotations

import html as stdhtml
import io
from typing import Any

import markdown as md_lib
from docx import Document
from htmldocx import HtmlToDocx
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from weasyprint import HTML as WeasyHTML


# ── PDF ─────────────────────────────────────────────────────────

def render_pdf(
    *,
    title: str,
    markdown_content: str | None = None,
    html_content: str | None = None,
) -> tuple[bytes, str]:
    """Render a PDF from markdown or raw HTML.

    Returns `(pdf_bytes, preview_html)`. The preview HTML is the same
    document used as input to WeasyPrint, so the inline preview is
    byte-for-byte the same layout users see when they open the PDF.
    """
    if html_content is not None:
        inner = html_content
    else:
        inner = md_lib.markdown(
            markdown_content or "",
            extensions=["extra", "sane_lists", "tables", "fenced_code"],
        )
    full_html = _wrap_html(title, inner)
    pdf_bytes = WeasyHTML(string=full_html).write_pdf()
    return pdf_bytes, full_html


# ── DOCX ────────────────────────────────────────────────────────

def render_docx(*, title: str, markdown_content: str) -> tuple[bytes, str]:
    """Render a Word document from markdown.

    Returns `(docx_bytes, preview_html)`. The preview HTML uses the same
    wrapper as the PDF path so users see consistent styling regardless
    of which format they generate.
    """
    inner = md_lib.markdown(
        markdown_content or "",
        extensions=["extra", "sane_lists", "tables", "fenced_code"],
    )
    full_html = _wrap_html(title, inner)

    doc = Document()
    parser = HtmlToDocx()
    parser.add_html_to_document(inner, doc)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), full_html


# ── XLSX ────────────────────────────────────────────────────────

def render_xlsx(
    *, title: str, sheets: list[dict[str, Any]]
) -> tuple[bytes, str]:
    """Render a workbook from a structured spec.

    Each sheet is `{"name": str, "headers": list[str], "rows": list[list]}`.
    Returns `(xlsx_bytes, preview_html)` where preview_html is a self-contained
    HTML document with one <table> per sheet and clickable sheet tabs.
    """
    if not sheets:
        raise ValueError("render_xlsx requires at least one sheet")

    wb = Workbook()
    # openpyxl auto-creates a default sheet; remove it and build ours.
    default = wb.active
    if default is not None:
        wb.remove(default)

    for s in sheets:
        sheet_name = str(s.get("name") or "Sheet")[:31]
        ws = wb.create_sheet(title=sheet_name)
        headers = s.get("headers") or []
        if headers:
            ws.append([_coerce_cell(h) for h in headers])
        for row in s.get("rows") or []:
            ws.append([_coerce_cell(v) for v in row])

    buf = io.BytesIO()
    wb.save(buf)
    binary = buf.getvalue()

    preview_html = _xlsx_preview_html(title, sheets)
    return binary, preview_html


# ── PPTX ────────────────────────────────────────────────────────

# Supported slide layouts. Keep the list small and explicit so the
# agent can reliably pick one and so the preview HTML can render a
# faithful thumbnail for each.
_PPTX_LAYOUTS = {"title", "section", "bullets", "content", "two_column"}


def render_pptx(
    *, title: str, slides: list[dict[str, Any]]
) -> tuple[bytes, str]:
    """Render a PowerPoint deck from a structured slide spec.

    Each slide is a dict with a `layout` field plus layout-specific
    fields. Supported layouts:

    - **title**: `{layout: "title", title: str, subtitle?: str}`
    - **section**: `{layout: "section", title: str, subtitle?: str}`
    - **bullets**: `{layout: "bullets", title: str, bullets: list[str]}`
    - **content**: `{layout: "content", title: str, body: str}`
    - **two_column**: `{layout: "two_column", title: str, left: list[str], right: list[str]}`

    Returns `(pptx_bytes, preview_html)`. The preview is a self-contained
    HTML document rendering each slide as a 16:9 card in a vertical
    scroll list — designed to look like slide thumbnails.
    """
    if not slides:
        raise ValueError("render_pptx requires at least one slide")

    prs = Presentation()
    # Use the default 10x7.5 inch slide size (4:3). That's what the
    # default template gives us; we could switch to 13.33x7.5 (16:9)
    # but 4:3 is more universally compatible for v1. The preview HTML
    # uses 16:9 cards for visual polish — that's a preview-only choice.

    for s in slides:
        layout = str(s.get("layout") or "bullets")
        slide_title = str(s.get("title") or "")

        if layout == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            _set_placeholder_text(slide, 0, slide_title)
            _set_placeholder_text(slide, 1, str(s.get("subtitle") or ""))

        elif layout == "section":
            # Section Header layout (index 2 in the default template)
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            _set_placeholder_text(slide, 0, slide_title)
            _set_placeholder_text(slide, 1, str(s.get("subtitle") or ""))

        elif layout == "bullets":
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            _set_placeholder_text(slide, 0, slide_title)
            bullets = s.get("bullets") or []
            _set_bullets(slide, 1, [str(b) for b in bullets])

        elif layout == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            _set_placeholder_text(slide, 0, slide_title)
            _set_placeholder_text(slide, 1, str(s.get("body") or ""))

        elif layout == "two_column":
            # Two Content layout (index 3 in the default template)
            slide = prs.slides.add_slide(prs.slide_layouts[3])
            _set_placeholder_text(slide, 0, slide_title)
            left = [str(x) for x in (s.get("left") or [])]
            right = [str(x) for x in (s.get("right") or [])]
            _set_bullets(slide, 1, left)
            _set_bullets(slide, 2, right)

        else:
            # Unknown layout → fall back to bullets so we never crash
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            _set_placeholder_text(slide, 0, slide_title or layout)
            _set_placeholder_text(slide, 1, str(s.get("body") or ""))

    buf = io.BytesIO()
    prs.save(buf)
    binary = buf.getvalue()
    preview_html = _pptx_preview_html(title, slides)
    return binary, preview_html


def _set_placeholder_text(slide, idx: int, text: str) -> None:
    """Safely set text on a placeholder by index — layouts may not
    have every placeholder, so we no-op if it's missing."""
    try:
        ph = slide.placeholders[idx]
    except (KeyError, IndexError):
        return
    if ph.has_text_frame:
        ph.text_frame.text = text
    else:
        ph.text = text


def _set_bullets(slide, idx: int, items: list[str]) -> None:
    """Set a placeholder's text frame to a bulleted list."""
    try:
        ph = slide.placeholders[idx]
    except (KeyError, IndexError):
        return
    if not ph.has_text_frame:
        return
    tf = ph.text_frame
    if not items:
        tf.text = ""
        return
    tf.text = items[0]
    for line in items[1:]:
        p = tf.add_paragraph()
        p.text = line


# ── Helpers ─────────────────────────────────────────────────────

def _coerce_cell(value: Any) -> Any:
    """Pass numbers/bools through; stringify everything else to avoid
    openpyxl exceptions on unexpected types (dicts, lists, etc.)."""
    if value is None:
        return None
    if isinstance(value, (int, float, bool)):
        return value
    return str(value)


def _wrap_html(title: str, body_html: str) -> str:
    """Wrap parsed markdown in a self-contained HTML document with
    print-friendly styling. Shared between WeasyPrint (PDF pipeline)
    and the inline preview so both look identical."""
    safe_title = stdhtml.escape(title)
    return f"""<!doctype html>
<html><head><meta charset="utf-8"><title>{safe_title}</title>
<style>
  @page {{ size: letter; margin: 0.75in; }}
  body {{
    font-family: "DejaVu Sans", "Liberation Sans", sans-serif;
    color: #1a1a1a;
    line-height: 1.5;
    max-width: 720px;
    margin: 2rem auto;
    padding: 0 1rem;
    background: #ffffff;
  }}
  h1 {{ font-size: 1.8rem; border-bottom: 2px solid #ccc; padding-bottom: .3rem; }}
  h2 {{ font-size: 1.4rem; margin-top: 1.6rem; }}
  h3 {{ font-size: 1.15rem; }}
  p, li {{ font-size: 0.95rem; }}
  code {{ background: #f4f4f4; padding: 0 .2em; border-radius: 3px;
          font-family: "DejaVu Sans Mono", monospace; }}
  pre {{ background: #f4f4f4; padding: .8rem; border-radius: 4px; overflow-x: auto; }}
  pre code {{ background: transparent; padding: 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; font-size: .9rem; }}
  th, td {{ border: 1px solid #ccc; padding: .4rem .6rem; text-align: left; }}
  th {{ background: #f0f0f0; }}
  blockquote {{ border-left: 3px solid #ccc; padding-left: 1rem; color: #555; }}
  a {{ color: #0066cc; }}
</style></head>
<body>{body_html}</body></html>"""


def _xlsx_preview_html(title: str, sheets: list[dict[str, Any]]) -> str:
    """Render an XLSX workbook spec as a tabbed, self-contained HTML document.

    Safe to embed inside an iframe via srcDoc with
    `sandbox="allow-scripts"` — the tab-switching helper is the only script.
    All cell values are html-escaped to prevent XSS even though the data
    originates from a trusted agent.
    """
    safe_title = stdhtml.escape(title)

    tabs: list[str] = []
    panels: list[str] = []
    for idx, s in enumerate(sheets):
        sheet_name = str(s.get("name") or f"Sheet{idx + 1}")
        safe_name = stdhtml.escape(sheet_name)
        active = "active" if idx == 0 else ""
        tabs.append(
            f'<button type="button" class="tab {active}" '
            f'data-idx="{idx}" onclick="show({idx})">{safe_name}</button>'
        )

        headers = s.get("headers") or []
        rows = s.get("rows") or []
        header_html = "".join(
            f"<th>{stdhtml.escape(str(h))}</th>" for h in headers
        )
        row_html_parts: list[str] = []
        for row in rows:
            cells = "".join(
                f"<td>{stdhtml.escape(str(v)) if v is not None else ''}</td>"
                for v in row
            )
            row_html_parts.append(f"<tr>{cells}</tr>")
        row_html = "".join(row_html_parts)

        display = "" if idx == 0 else "display:none"
        panels.append(
            f'<div class="panel" id="p{idx}" style="{display}">'
            f"<table><thead><tr>{header_html}</tr></thead>"
            f"<tbody>{row_html}</tbody></table></div>"
        )

    tabs_html = "".join(tabs)
    panels_html = "".join(panels)

    return f"""<!doctype html>
<html><head><meta charset="utf-8"><title>{safe_title}</title>
<style>
  body {{
    font-family: "DejaVu Sans", "Liberation Sans", sans-serif;
    color: #1a1a1a;
    margin: 0;
    background: #ffffff;
  }}
  .tabs {{
    display: flex;
    gap: 2px;
    background: #f0f0f0;
    border-bottom: 1px solid #ccc;
    padding: 0.5rem 0.75rem 0;
    overflow-x: auto;
  }}
  .tab {{
    background: #e6e6e6;
    border: 1px solid #ccc;
    border-bottom: none;
    border-radius: 4px 4px 0 0;
    padding: 0.4rem 0.9rem;
    font-size: 0.85rem;
    cursor: pointer;
    color: #444;
    font-family: inherit;
  }}
  .tab:hover {{ background: #f4f4f4; }}
  .tab.active {{
    background: #ffffff;
    color: #1a1a1a;
    font-weight: 600;
    position: relative;
    z-index: 1;
  }}
  .panel {{
    padding: 1rem;
    overflow-x: auto;
  }}
  table {{
    border-collapse: collapse;
    width: 100%;
    font-size: 0.85rem;
  }}
  th, td {{
    border: 1px solid #ddd;
    padding: 0.35rem 0.55rem;
    text-align: left;
    white-space: nowrap;
  }}
  thead th {{
    background: #f0f0f0;
    font-weight: 600;
    position: sticky;
    top: 0;
  }}
  tbody tr:nth-child(even) {{ background: #fafafa; }}
</style></head>
<body>
<div class="tabs">{tabs_html}</div>
{panels_html}
<script>
  function show(i) {{
    document.querySelectorAll(".panel").forEach(function (p, j) {{
      p.style.display = j === i ? "" : "none";
    }});
    document.querySelectorAll(".tab").forEach(function (t, j) {{
      if (j === i) t.classList.add("active"); else t.classList.remove("active");
    }});
  }}
</script>
</body></html>"""


def _pptx_preview_html(title: str, slides: list[dict[str, Any]]) -> str:
    """Render a PPTX deck spec as a self-contained HTML document.

    Each slide becomes a 16:9 "thumbnail" card in a vertical scroll list.
    All user data is html.escape()'d — safe to sandbox in an iframe
    via srcDoc with `sandbox="allow-scripts"`.
    """
    safe_title = stdhtml.escape(title)

    cards: list[str] = []
    for idx, s in enumerate(slides):
        layout = str(s.get("layout") or "bullets")
        slide_title = stdhtml.escape(str(s.get("title") or ""))
        slide_num = idx + 1

        if layout == "title":
            subtitle = stdhtml.escape(str(s.get("subtitle") or ""))
            body_html = (
                f'<div class="slide-title-body">'
                f'<h1 class="title-hero">{slide_title}</h1>'
                f'{f"<p class=\'title-sub\'>{subtitle}</p>" if subtitle else ""}'
                f"</div>"
            )
        elif layout == "section":
            subtitle = stdhtml.escape(str(s.get("subtitle") or ""))
            body_html = (
                f'<div class="slide-section-body">'
                f'<div class="section-label">Section</div>'
                f'<h2 class="section-title">{slide_title}</h2>'
                f'{f"<p class=\'section-sub\'>{subtitle}</p>" if subtitle else ""}'
                f"</div>"
            )
        elif layout == "bullets":
            bullets = s.get("bullets") or []
            bullet_html = "".join(
                f"<li>{stdhtml.escape(str(b))}</li>" for b in bullets
            )
            body_html = (
                f'<h3 class="slide-heading">{slide_title}</h3>'
                f'<ul class="slide-bullets">{bullet_html}</ul>'
            )
        elif layout == "content":
            body_text = stdhtml.escape(str(s.get("body") or ""))
            body_html = (
                f'<h3 class="slide-heading">{slide_title}</h3>'
                f'<p class="slide-body">{body_text}</p>'
            )
        elif layout == "two_column":
            left = s.get("left") or []
            right = s.get("right") or []
            left_html = "".join(
                f"<li>{stdhtml.escape(str(x))}</li>" for x in left
            )
            right_html = "".join(
                f"<li>{stdhtml.escape(str(x))}</li>" for x in right
            )
            body_html = (
                f'<h3 class="slide-heading">{slide_title}</h3>'
                f'<div class="slide-two-col">'
                f'<ul class="slide-bullets">{left_html}</ul>'
                f'<ul class="slide-bullets">{right_html}</ul>'
                f"</div>"
            )
        else:
            body_text = stdhtml.escape(str(s.get("body") or ""))
            body_html = (
                f'<h3 class="slide-heading">{slide_title or layout}</h3>'
                f'<p class="slide-body">{body_text}</p>'
            )

        cards.append(
            f'<div class="slide-card">'
            f'<div class="slide-number">{slide_num}</div>'
            f'<div class="slide-content">{body_html}</div>'
            f"</div>"
        )

    cards_html = "".join(cards)
    return f"""<!doctype html>
<html><head><meta charset="utf-8"><title>{safe_title}</title>
<style>
  body {{
    font-family: "DejaVu Sans", "Liberation Sans", sans-serif;
    color: #1a1a1a;
    margin: 0;
    padding: 1rem;
    background: #f4f4f4;
  }}
  .slide-card {{
    position: relative;
    width: 100%;
    max-width: 800px;
    margin: 0 auto 1rem;
    aspect-ratio: 16 / 9;
    background: #ffffff;
    border: 1px solid #d4d4d4;
    border-radius: 6px;
    box-shadow: 0 2px 6px rgba(0, 0, 0, 0.08);
    padding: 2.4rem 3rem;
    box-sizing: border-box;
    overflow: hidden;
  }}
  .slide-number {{
    position: absolute;
    top: 0.6rem;
    right: 0.9rem;
    font-size: 0.72rem;
    color: #888;
    font-weight: 500;
  }}
  .slide-content {{ height: 100%; display: flex; flex-direction: column; }}
  .slide-heading {{
    font-size: 1.35rem;
    font-weight: 600;
    color: #1a1a1a;
    margin: 0 0 1.1rem 0;
    padding-bottom: 0.5rem;
    border-bottom: 2px solid #e6e6e6;
  }}
  .slide-bullets {{
    margin: 0;
    padding-left: 1.3rem;
    font-size: 0.95rem;
    line-height: 1.65;
    color: #333;
  }}
  .slide-bullets li {{ margin-bottom: 0.35rem; }}
  .slide-body {{
    font-size: 0.95rem;
    line-height: 1.6;
    color: #333;
    margin: 0;
    white-space: pre-wrap;
  }}
  .slide-two-col {{ display: flex; gap: 2rem; flex: 1; }}
  .slide-two-col > ul {{ flex: 1; }}
  .slide-title-body {{
    display: flex;
    flex-direction: column;
    justify-content: center;
    align-items: center;
    text-align: center;
    height: 100%;
  }}
  .title-hero {{
    font-size: 2.2rem;
    font-weight: 700;
    margin: 0 0 0.6rem 0;
    color: #1a1a1a;
  }}
  .title-sub {{ font-size: 1.1rem; color: #666; margin: 0; }}
  .slide-section-body {{
    display: flex;
    flex-direction: column;
    justify-content: center;
    height: 100%;
    border-left: 6px solid #1a1a1a;
    padding-left: 1.5rem;
  }}
  .section-label {{
    font-size: 0.7rem;
    text-transform: uppercase;
    letter-spacing: 0.15em;
    color: #888;
    margin-bottom: 0.4rem;
  }}
  .section-title {{
    font-size: 1.9rem;
    font-weight: 700;
    color: #1a1a1a;
    margin: 0 0 0.4rem 0;
  }}
  .section-sub {{ font-size: 0.95rem; color: #555; margin: 0; }}
</style></head>
<body>{cards_html}</body></html>"""
